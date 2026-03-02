#!/usr/bin/env python3
"""Create 5 PowerPoint slides summarizing the Citicorp Center simulation.

Slide 1: Building overview + 3D structural system figure
Slide 2: FEA & Connection analysis with figures
Slide 3: 16-Scenario Risk Matrix (Monte Carlo heatmaps)
Slide 4: 1970s vs Modern wind analysis comparison
Slide 5: Validation, provenance, and key findings

Uses python-pptx. Embeds MATLAB-generated figures from figures/ folder.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os
import json

# ============================================================
# STYLE CONSTANTS
# ============================================================
BG_COLOR = RGBColor(0x1A, 0x1A, 0x2E)
TITLE_COLOR = RGBColor(0xFF, 0xFF, 0xFF)
SUBTITLE_COLOR = RGBColor(0xAA, 0xCC, 0xFF)
BODY_COLOR = RGBColor(0xDD, 0xDD, 0xDD)
ACCENT_GREEN = RGBColor(0x33, 0xCC, 0x33)
ACCENT_RED = RGBColor(0xFF, 0x44, 0x44)
ACCENT_GOLD = RGBColor(0xFF, 0xCC, 0x33)
ACCENT_CYAN = RGBColor(0x33, 0xCC, 0xFF)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

FIGURES_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'figures')


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=14,
                 color=BODY_COLOR, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_rich_text_box(slide, left, top, width, height):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    return tf


def add_para(tf, text, font_size=12, color=BODY_COLOR, bold=False,
             alignment=PP_ALIGN.LEFT, space_before=Pt(2), space_after=Pt(2),
             font_name='Calibri'):
    if len(tf.paragraphs) == 1 and tf.paragraphs[0].text == '':
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    return p


def add_table(slide, left, top, width, height, rows, cols):
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    return table_shape.table


def style_cell(cell, text, font_size=10, color=BODY_COLOR, bold=False,
               fill_color=None, alignment=PP_ALIGN.CENTER, font_name='Calibri'):
    cell.text = text
    for paragraph in cell.text_frame.paragraphs:
        paragraph.font.size = Pt(font_size)
        paragraph.font.color.rgb = color
        paragraph.font.bold = bold
        paragraph.font.name = font_name
        paragraph.alignment = alignment
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if fill_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill_color


def add_figure(slide, filename, left, top, width, height):
    """Embed a figure from the figures/ folder. Returns True if found."""
    path = os.path.join(FIGURES_DIR, filename)
    if os.path.exists(path):
        slide.shapes.add_picture(path, left, top, width, height)
        return True
    return False


# ============================================================
# SLIDE 1: Building Overview + 3D Figure
# ============================================================
def create_slide_1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_COLOR)

    add_text_box(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.6),
                 'Citicorp Center Wind Loading Crisis Simulation',
                 font_size=28, color=TITLE_COLOR, bold=True,
                 alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0.5), Inches(0.7), Inches(12), Inches(0.35),
                 'OR 750: Reliability, Safety, and Risk  |  George Mason University  |  ASCE 7-22 / AISC 360-22',
                 font_size=13, color=SUBTITLE_COLOR, alignment=PP_ALIGN.CENTER)

    # Left column: Building parameters + crisis narrative
    tf1 = add_rich_text_box(slide, Inches(0.3), Inches(1.2), Inches(4.8), Inches(6.0))
    add_para(tf1, 'Building Parameters', 15, ACCENT_CYAN, True)

    params = [
        ('Height', '915 ft (59 stories)'),
        ('Plan', '157 x 157 ft (square)'),
        ('Stilts', '114 ft, columns at mid-face'),
        ('Cantilever', '72 ft at each corner'),
        ('Bracing', '6 tiers of chevron V (48 braces)'),
        ('Transfer', '2-story W-pattern truss'),
        ('TMD', '400 tons on floor 63'),
        ('Period', '6.5 sec fundamental'),
    ]
    for label, value in params:
        add_para(tf1, f'  {label}: {value}', 10, BODY_COLOR)

    add_para(tf1, '', 6, BODY_COLOR)
    add_para(tf1, 'The 1978 Crisis', 15, ACCENT_RED, True)
    crisis = [
        'Original design: CJP groove welds',
        'Bethlehem Steel substituted A325 bolts',
        'Bolts sized for perpendicular wind ONLY',
        'LeMessurier: quartering never checked',
        'Emergency repair during Hurricane Ella',
        'Secret repair, 2" welded cover plates',
    ]
    for c in crisis:
        add_para(tf1, f'  \u2022 {c}', 10, BODY_COLOR)

    add_para(tf1, '', 6, BODY_COLOR)
    add_para(tf1, 'Simulation Modules', 13, ACCENT_CYAN, True)
    modules = [
        '1: 3D Visualization', '2: ASCE 7-22 Wind',
        '3: 3D FEA', '4: AISC 360-22 Connections',
        '5: Monte Carlo (16 scenarios)',
        '6: Validation', '10: 1978 vs Modern',
        '11: Provenance',
    ]
    for m in modules:
        add_para(tf1, f'  {m}', 9, RGBColor(0xAA, 0xAA, 0xAA))

    # Right side: 3D building figure
    add_figure(slide, '01_building_3d.png',
               Inches(5.3), Inches(1.1), Inches(7.8), Inches(5.8))

    # Footer
    add_text_box(slide, Inches(0.5), Inches(7.1), Inches(12), Inches(0.25),
                 'Sources: NIST (2021), Morgenstern (1995), ASCE 7-22, AISC 360-22  |  Base MATLAB only',
                 font_size=8, color=RGBColor(0x66, 0x66, 0x66),
                 alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 2: FEA & Connection Analysis with Figures
# ============================================================
def create_slide_2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_COLOR)

    add_text_box(slide, Inches(0.3), Inches(0.15), Inches(12.5), Inches(0.55),
                 'Structural Analysis: FEA & Connection Capacity',
                 font_size=26, color=TITLE_COLOR, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # FEA figure (left half)
    add_figure(slide, '03_fea_chevron.png',
               Inches(0.2), Inches(0.85), Inches(6.4), Inches(4.2))

    # Connection figure (right half)
    add_figure(slide, '04_connections.png',
               Inches(6.7), Inches(0.85), Inches(6.4), Inches(4.2))

    # Key metrics bar at bottom
    tf = add_rich_text_box(slide, Inches(0.3), Inches(5.2), Inches(12.5), Inches(2.0))
    add_para(tf, 'Key Results', 14, ACCENT_CYAN, True)

    metrics = [
        'FEA: 80 nodes, 264 elements, equilibrium error = 4.8e-15 (PASS)',
        'Max displacement: 12.3 in (perp.), 4.6 in (quartering) | Drift: H/890',
        'FEA quartering amplification: 0.52x (face wind governs in modern analysis)',
        'Historical quartering amp: 1.40x (LeMessurier 1978, static superposition)',
        'D/C (A325 Bolted, J3-3a): Perp = 0.54 (OK), Quart = 0.24 (OK under modern loads)',
        'D/C (CJP Weld): Perp = 0.31, Quart = 0.16  |  Repaired: Perp = 0.27, Quart = 0.14',
    ]
    for m in metrics:
        add_para(tf, f'  \u2022 {m}', 10, BODY_COLOR)


# ============================================================
# SLIDE 3: 16-Scenario Risk Matrix
# ============================================================
def create_slide_3(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_COLOR)

    add_text_box(slide, Inches(0.3), Inches(0.15), Inches(12.5), Inches(0.55),
                 '16-Scenario Monte Carlo Reliability (2\u00d72\u00d72\u00d72 Factorial)',
                 font_size=26, color=TITLE_COLOR, bold=True,
                 alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0.3), Inches(0.65), Inches(12.5), Inches(0.3),
                 'Methodology (1970s/Modern) \u00d7 Connection (Weld/Bolt) \u00d7 Direction (Perp/Quar) \u00d7 TMD (ON/OFF)',
                 font_size=12, color=ACCENT_GOLD, alignment=PP_ALIGN.CENTER)

    # Main Monte Carlo figure
    add_figure(slide, '05_monte_carlo.png',
               Inches(0.2), Inches(1.0), Inches(13.0), Inches(6.0))

    # If JSON results exist, add a small summary overlay
    json_path = os.path.join(FIGURES_DIR, 'mc_results_16.json')
    if os.path.exists(json_path):
        try:
            with open(json_path, 'r') as f:
                mc = json.load(f)
            # Small text box with critical scenarios
            tf = add_rich_text_box(slide, Inches(0.3), Inches(7.05), Inches(12.5), Inches(0.35))
            T = mc.get('T_return', [])
            if len(T) >= 16:
                txt = (f'Critical: Sc.8 (70s-Bolt-Q-off) T={T[7]:.0f}yr | '
                       f'Sc.7 (70s-Bolt-Q-on) T={T[6]:.0f}yr | '
                       f'Sc.16 (Mod-Bolt-Q-off) T={T[15]:.0f}yr | '
                       f'Validation: Sc.8 vs hist T=16yr, Sc.7 vs hist T=55yr')
                add_para(tf, txt, 8, RGBColor(0x88, 0x88, 0x88),
                         alignment=PP_ALIGN.CENTER)
        except (json.JSONDecodeError, KeyError):
            pass


# ============================================================
# SLIDE 4: 1970s vs Modern Comparison
# ============================================================
def create_slide_4(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_COLOR)

    add_text_box(slide, Inches(0.3), Inches(0.15), Inches(12.5), Inches(0.55),
                 '1970s Analysis vs Modern Reassessment',
                 font_size=26, color=TITLE_COLOR, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # 1978 comparison figure
    add_figure(slide, '10_1978_comparison.png',
               Inches(0.1), Inches(0.8), Inches(8.5), Inches(5.5))

    # Right side: narrative explanation
    tf = add_rich_text_box(slide, Inches(8.8), Inches(0.9), Inches(4.3), Inches(6.0))

    add_para(tf, 'What LeMessurier Believed', 14, ACCENT_RED, True)
    add_para(tf, '(1978 static superposition)', 10, RGBColor(0xFF, 0x88, 0x66))
    beliefs = [
        'Quartering = 1.40\u00d7 perpendicular',
        'Full pressure on BOTH faces',
        'Treats loads as perfectly correlated',
        'Bolt D/C >> 1.0 under quartering',
        'T \u2248 16 yr without TMD (crisis)',
        'T \u2248 55 yr with TMD (marginal)',
    ]
    for b in beliefs:
        add_para(tf, f'  \u2022 {b}', 10, BODY_COLOR)

    add_para(tf, '', 6, BODY_COLOR)
    add_para(tf, 'What Modern Analysis Shows', 14, ACCENT_GREEN, True)
    add_para(tf, '(ASCE 7-22 + NIST 2021)', 10, RGBColor(0x88, 0xFF, 0x88))
    modern = [
        'Quartering = 0.52\u00d7 perpendicular',
        'Face winds actually govern',
        'Kd = 0.85 reduces all loads by 15%',
        'Static superposition was conservative',
        'But bolts were still under-designed!',
        'The real problem: welds \u2192 bolts swap',
    ]
    for m in modern:
        add_para(tf, f'  \u2022 {m}', 10, BODY_COLOR)

    add_para(tf, '', 6, BODY_COLOR)
    add_para(tf, 'Key Takeaway', 13, ACCENT_GOLD, True)
    add_para(tf, ('LeMessurier was right to be alarmed, '
                  'even if the quartering assumption was '
                  'conservative. The bolt substitution was '
                  'the root cause regardless of methodology.'),
             10, BODY_COLOR)


# ============================================================
# SLIDE 5: Validation & Key Findings
# ============================================================
def create_slide_5(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_COLOR)

    add_text_box(slide, Inches(0.3), Inches(0.15), Inches(12.5), Inches(0.55),
                 'Validation, Provenance & Key Findings',
                 font_size=26, color=TITLE_COLOR, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # Validation figure (left)
    add_figure(slide, '06_validation.png',
               Inches(0.1), Inches(0.8), Inches(6.4), Inches(4.0))

    # Wind pressure figure (right)
    add_figure(slide, '02_wind_pressure.png',
               Inches(6.7), Inches(0.8), Inches(6.4), Inches(4.0))

    # Bottom: Key findings and provenance summary
    tf = add_rich_text_box(slide, Inches(0.3), Inches(5.0), Inches(6.0), Inches(2.3))
    add_para(tf, 'Key Findings for OR 750', 14, ACCENT_CYAN, True)

    findings = [
        ('Risk Quantification:', 'Full factorial (16 scenarios) reveals '
         'that methodology choice dramatically affects perceived risk'),
        ('Actual vs Perceived:', '1970s analysis overestimates quartering risk; '
         'modern analysis shows face winds govern'),
        ('Root Cause:', 'Weld-to-bolt substitution created under-designed '
         'connections regardless of wind direction methodology'),
        ('TMD Importance:', 'TMD provides significant risk reduction '
         'across all scenarios (damping 1% \u2192 8%)'),
        ('No Circular Reasoning:', 'Gumbel from ASCE 7 wind map; failure '
         'thresholds from FEA+capacity; historical values for validation only'),
    ]
    for title, detail in findings:
        add_para(tf, f'  {title}', 10, ACCENT_GOLD, True, space_before=Pt(4))
        add_para(tf, f'    {detail}', 9, BODY_COLOR, space_before=Pt(0))

    # Data provenance summary
    tf2 = add_rich_text_box(slide, Inches(6.5), Inches(5.0), Inches(6.5), Inches(2.3))
    add_para(tf2, 'Data Provenance Categories', 14, ACCENT_CYAN, True)

    cats = [
        ('\u2588 SOURCED', 'Published standards/references (ASCE 7-22, AISC 360-22)',
         RGBColor(0x33, 0xCC, 0x33)),
        ('\u2588 STANDARD', 'Code-prescribed values (Kd, Cp, bolt capacities)',
         RGBColor(0x33, 0x99, 0xFF)),
        ('\u2588 ASSUMED', 'Engineering estimates (sections, damping ratios)',
         RGBColor(0xFF, 0xCC, 0x33)),
        ('\u2588 DERIVED', 'Computed in this simulation (quartering amp, T_return)',
         RGBColor(0xFF, 0x66, 0xCC)),
        ('\u2588 VALIDATED', 'Compared to independent data (LeMessurier, NIST)',
         RGBColor(0xFF, 0xFF, 0xFF)),
    ]
    for symbol, desc, color in cats:
        p = add_para(tf2, '', 10, BODY_COLOR, space_before=Pt(3))
        run1 = p.add_run()
        run1.text = f'{symbol} '
        run1.font.size = Pt(10)
        run1.font.color.rgb = color
        run1.font.bold = True
        run1.font.name = 'Calibri'
        run2 = p.add_run()
        run2.text = desc
        run2.font.size = Pt(9)
        run2.font.color.rgb = BODY_COLOR
        run2.font.name = 'Calibri'

    # Footer
    add_text_box(slide, Inches(0.3), Inches(7.1), Inches(12.5), Inches(0.25),
                 'Refs: ASCE 7-22, AISC 360-22, NIST (2021), Morgenstern (1995)  |  Generated with Claude Code',
                 font_size=8, color=RGBColor(0x66, 0x66, 0x66),
                 alignment=PP_ALIGN.CENTER)


# ============================================================
# MAIN
# ============================================================
def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    create_slide_1(prs)
    create_slide_2(prs)
    create_slide_3(prs)
    create_slide_4(prs)
    create_slide_5(prs)

    out_dir = os.path.dirname(os.path.abspath(__file__))
    out_path = os.path.join(out_dir, 'citicorp_16scenario_summary.pptx')
    prs.save(out_path)
    print(f'Saved: {out_path}')
    print(f'  5 slides with embedded figures, widescreen (13.33" x 7.5")')


if __name__ == '__main__':
    main()
