#!/usr/bin/env python3
"""Create PowerPoint with methodology + ParaView old/new comparison views.

Regenerates cityscape_comparison.pptx with updated v3 screenshots and
full methodology slides.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
PARENT_DIR = os.path.dirname(SCRIPT_DIR)
OUT_PATH = os.path.join(PARENT_DIR, 'cityscape_comparison_v2.pptx')

NAVY = RGBColor(0x1B, 0x3A, 0x5C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x55, 0x55, 0x55)
DARK = RGBColor(0x33, 0x33, 0x33)
LIGHT_BLUE = RGBColor(0xBB, 0xCC, 0xDD)

# v3 screenshots (latest, with rotated tower + spatial dedup)
COMPARISONS = [
    ('v3_old_aerial.png',   'v3_new_aerial.png',   'Aerial View'),
    ('v3_old_quarter.png',  'v3_new_quarter.png',   '3/4 View'),
    ('v3_old_street.png',   'v3_new_street.png',    'Street Level'),
    ('v3_old_top_down.png', 'v3_new_top_down.png',  'Plan View (Top Down)'),
    ('v3_old_stilt.png',    'v3_new_stilt.png',     'Stilt Zone Detail'),
]


def textbox(slide, x, y, w, h, text, size=14, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    """Add a text box and return the paragraph."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return tf


def bullet_list(slide, x, y, w, h, items, size=11, color=DARK):
    """Add a multi-line bullet list."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, text in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(4)
    return tf


def slide_title(slide, text):
    textbox(slide, 0.3, 0.15, 9.4, 0.5, text, size=22, bold=True, color=NAVY, align=PP_ALIGN.CENTER)


# ========== SLIDE BUILDERS ==========

def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY
    textbox(slide, 0.5, 1.2, 9, 1.5,
            "Citicorp Center CFD\nCityscape Geometry Comparison",
            size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    textbox(slide, 0.5, 3.2, 9, 1,
            "Old (Socrata LoD1 Extrusions) vs New (CityGML LoD2 Hybrid)",
            size=22, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
    textbox(slide, 0.5, 4.5, 9, 0.5,
            "OR 750 -- Reliability, Safety, and Risk",
            size=14, color=RGBColor(0x88, 0x99, 0xAA), align=PP_ALIGN.CENTER)


def add_methodology_overview(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_title(slide, "Methodology: Hybrid STL Pipeline")
    bullet_list(slide, 0.5, 0.7, 9, 4.5, [
        "1. CITICORP TOWER + STILTS (hand-crafted, photo-refined)",
        "     Tower: 47.85m square plan, rotated 28.7 deg to match Manhattan grid",
        "     4 rectangular stilts at face midpoints (7.27m x 6.55m, flush with face)",
        "     Octagonal center core (21.09m x 21.09m, 3.27m chamfers)",
        "",
        "2. CITYGML SURROUNDINGS (primary source -- 2014 aerial photogrammetry)",
        "     NYC 3D Building Model: LoD2 with WallSurface, RoofSurface, GroundSurface",
        "     Stream ~764 MB GML file (DA12), extract buildings in CFD domain",
        "     EPSG:2263 (NY State Plane) -> local meters centered at Citicorp",
        "     Result: 599 buildings, 55,382 triangles -- actual roof shapes + setbacks",
        "",
        "3. SOCRATA API FALLBACK (for post-2014 construction)",
        "     NYC Building Footprints API: continuously updated, GeoJSON with height",
        "     Deduplicate by BIN match (499 caught) + spatial proximity <15m (84 caught)",
        "     Extrude footprint polygon to roof height (LoD1, flat-topped)",
        "     Result: 123 new buildings, 4,828 triangles",
    ], size=10)


def add_data_sources(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_title(slide, "Data Sources")

    # CityGML box
    textbox(slide, 0.4, 0.7, 4.3, 0.35, "NYC 3D Building Model (CityGML)", size=14, bold=True, color=NAVY)
    bullet_list(slide, 0.6, 1.05, 4.1, 2.5, [
        "2014 aerial photogrammetry survey",
        "~1 million buildings across 5 boroughs",
        "LoD2: actual wall, roof, ground surfaces",
        "12 DA district files (~6 GB total)",
        "EPSG:2263 coordinate system (US Survey Feet)",
        "Source: georocket/new-york-city-model-enhanced",
    ], size=10)

    # Socrata box
    textbox(slide, 5.2, 0.7, 4.3, 0.35, "Socrata Building Footprints API", size=14, bold=True, color=NAVY)
    bullet_list(slide, 5.4, 1.05, 4.1, 2.5, [
        "Continuously updated (includes 2015+ construction)",
        "2D footprint polygon + height_roof attribute",
        "BIN (Building Identification Number) for matching",
        "GeoJSON format, WGS84 coordinates",
        "Extruded to flat-topped LoD1 boxes",
        "Covers buildings missing from 2014 aerial survey",
    ], size=10)

    # Dedup box
    textbox(slide, 0.4, 3.5, 9.2, 0.35, "Deduplication Strategy", size=14, bold=True, color=NAVY)
    bullet_list(slide, 0.6, 3.85, 9, 1.5, [
        "Step 1: BIN matching -- if Socrata BIN exists in CityGML set, skip (caught 499)",
        "Step 2: Spatial proximity -- if Socrata centroid within 15m of CityGML centroid, skip (caught 84)",
        "Step 3: Remainder kept as fallback -- genuinely new post-2014 buildings (123 added)",
    ], size=10)


def add_stilt_methodology(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_title(slide, "Citicorp Stilt Geometry: Photo-Refined Dimensions")

    bullet_list(slide, 0.4, 0.7, 5, 4.5, [
        "Reference: street-level photo of Citicorp Center base",
        "",
        "Facade panel counting:",
        "  32 full + 2 half panels = 33 panel-widths across 157 ft",
        "  Panel = 48m / 33 = 1.4545 m",
        "",
        "Stilt columns (rectangular, NOT square):",
        "  Face width = 5 panels = 7.27 m",
        "  Depth = 4.5 panels = 6.55 m",
        "  Outer edge flush with tower face",
        "",
        "Center core (octagonal cross-section):",
        "  Bounding box = 14.5 panels = 21.09 m (square)",
        "  Flat edges = 10 panels = 14.55 m (each face)",
        "  Chamfer = 3.27 m corner cuts",
        "",
        "Tower rotation = 28.7 deg east of true north",
        "  Measured from CityGML LoD2 wall edge coordinates",
        "  Aligns with Manhattan street grid",
    ], size=10)

    # Octagon diagram placeholder
    textbox(slide, 5.5, 1.0, 4, 0.3, "Center Core Cross-Section", size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    bullet_list(slide, 5.8, 1.4, 3.5, 3, [
        "       10 panels",
        "    +-----------+",
        "   /             \\",
        "  |               |  14.5",
        "  |    21.09m     |  panels",
        "  |               |",
        "   \\             /",
        "    +-----------+",
        "       10 panels",
    ], size=9, color=NAVY)


def add_comparison_slide(prs, old_img, new_img, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_title(slide, title)

    img_w = Inches(4.7)
    img_h = Inches(2.65)
    y_pos = Inches(0.75)

    old_path = os.path.join(PARENT_DIR, old_img)
    new_path = os.path.join(PARENT_DIR, new_img)

    if os.path.exists(old_path):
        slide.shapes.add_picture(old_path, Inches(0.15), y_pos, img_w, img_h)
    if os.path.exists(new_path):
        slide.shapes.add_picture(new_path, Inches(5.15), y_pos, img_w, img_h)

    for label, x in [("OLD: Socrata LoD1 (706 buildings, 28k tris)", 0.15),
                      ("NEW: CityGML Hybrid (722 buildings, 60k tris)", 5.15)]:
        textbox(slide, x, 3.45, 4.7, 0.3, label, size=11, bold=True, color=DARK, align=PP_ALIGN.CENTER)

    bullet_list(slide, 0.3, 3.85, 9.4, 1.5, [
        "Red = Citicorp tower (rotated 28.7 deg to match grid)  |  Blue = Stilts + center core",
        "Gray = Socrata LoD1 surroundings  |  Green = CityGML LoD2 surroundings (actual roof shapes)",
    ], size=9, color=GRAY)


def add_numbers_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_title(slide, "Results: Old vs New Geometry")

    # Table-like layout
    headers = ["", "Old (Socrata LoD1)", "New (CityGML Hybrid)"]
    rows = [
        ["Surroundings source", "Socrata API only", "CityGML primary + Socrata fallback"],
        ["Building count", "706", "722 (599 CityGML + 123 Socrata)"],
        ["Surroundings tris", "28,376", "60,210"],
        ["Tower triangles", "12 (axis-aligned box)", "12 (rotated box) or 124 (LoD2)"],
        ["Stilt triangles", "48 (4 square columns)", "76 (4 rect stilts + octagonal core)"],
        ["Roof detail", "All flat-topped", "Actual roof shapes from photogrammetry"],
        ["Building setbacks", "None (extruded footprint)", "Wall setbacks at multiple Z levels"],
        ["Tower rotation", "Axis-aligned (wrong)", "28.7 deg (matches street grid)"],
        ["Stilt shape", "Square 7.32m", "Rectangular 7.27m x 6.55m"],
        ["Center core", "Not modeled", "21.09m octagon with chamfers"],
    ]

    y = 0.7
    # Headers
    for j, h in enumerate(headers):
        x = 0.3 + j * 3.1
        textbox(slide, x, y, 3.0, 0.3, h, size=11, bold=True, color=WHITE if j > 0 else NAVY,
                align=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)
        if j > 0:
            # Color header background
            tb = slide.shapes[-1]
            tb.fill.solid()
            tb.fill.fore_color.rgb = NAVY if j == 1 else RGBColor(0x2E, 0x7D, 0x32)

    y += 0.35
    for row in rows:
        for j, cell in enumerate(row):
            x = 0.3 + j * 3.1
            c = DARK if j == 0 else GRAY
            textbox(slide, x, y, 3.0, 0.22, cell, size=9, color=c,
                    align=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)
        y += 0.25


def add_summary_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY
    textbox(slide, 0.5, 1.5, 9, 1,
            "Key Improvements", size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    bullet_list(slide, 1, 2.8, 8, 2.5, [
        "CityGML LoD2 provides actual building surfaces (roof shapes, wall setbacks)",
        "Spatial deduplication eliminates 84 additional overlapping buildings",
        "Tower rotated 28.7 deg to match Manhattan street grid",
        "Photo-refined stilt dimensions: rectangular columns + octagonal center core",
        "2x triangle count (60k vs 28k) for more accurate flow modeling",
    ], size=14, color=WHITE)


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    add_title_slide(prs)           # 1
    add_methodology_overview(prs)  # 2
    add_data_sources(prs)          # 3
    add_stilt_methodology(prs)     # 4

    for old_img, new_img, title in COMPARISONS:
        add_comparison_slide(prs, old_img, new_img, title)  # 5-9

    add_numbers_slide(prs)         # 10
    add_summary_slide(prs)         # 11

    prs.save(OUT_PATH)
    print(f"Saved: {OUT_PATH} ({len(prs.slides)} slides)")


if __name__ == '__main__':
    main()
