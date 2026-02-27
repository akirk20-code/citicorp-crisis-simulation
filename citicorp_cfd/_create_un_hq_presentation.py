#!/usr/bin/env python3
"""Create PowerPoint for UN HQ CityGML extraction: methodology + views.

Generates un_hq_presentation.pptx with methodology, reproduction steps,
parameters, building inventory, and ParaView screenshots.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
PARENT_DIR = os.path.dirname(SCRIPT_DIR)
OUT_PATH = os.path.join(PARENT_DIR, 'un_hq_presentation.pptx')

NAVY = RGBColor(0x1B, 0x3A, 0x5C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x55, 0x55, 0x55)
DARK = RGBColor(0x33, 0x33, 0x33)
LIGHT_BLUE = RGBColor(0xBB, 0xCC, 0xDD)
GREEN = RGBColor(0x2E, 0x7D, 0x32)

SCREENSHOTS = [
    ('un_hq_aerial.png',      'Aerial View'),
    ('un_hq_quarter.png',     '3/4 View'),
    ('un_hq_street.png',      'Street Level'),
    ('un_hq_top_down.png',    'Plan View (Top Down)'),
    ('un_hq_east_river.png',  'East River Approach'),
]


def textbox(slide, x, y, w, h, text, size=14, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return tf


def bullet_list(slide, x, y, w, h, items, size=11, color=DARK):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, text in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(4)
    return tf


def slide_title(slide, text):
    textbox(slide, 0.3, 0.15, 9.4, 0.5, text, size=22, bold=True, color=NAVY, align=PP_ALIGN.CENTER)


# ========== SLIDE BUILDERS ==========

def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY
    textbox(slide, 0.5, 1.0, 9, 1.5,
            "UN Headquarters Area\nCityGML Hybrid Mesh Extraction",
            size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    textbox(slide, 0.5, 3.0, 9, 1,
            "CityGML LoD2 (primary) + Socrata API (fallback)",
            size=22, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
    textbox(slide, 0.5, 4.2, 9, 0.8,
            "Center for Computational Fluid Dynamics\n"
            "George Mason University",
            size=14, color=RGBColor(0x88, 0x99, 0xAA), align=PP_ALIGN.CENTER)


def add_parameters(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_title(slide, "Extraction Parameters")

    # Left column: location
    textbox(slide, 0.4, 0.7, 4.5, 0.35, "Location & Domain", size=14, bold=True, color=NAVY)
    bullet_list(slide, 0.6, 1.05, 4.3, 2.5, [
        "Center: 40.7489° N, 73.9680° W",
        "  (UN Secretariat Building centroid)",
        "Domain radius: 300 m (600 m × 600 m square)",
        "EPSG:2263 (NY State Plane): X=[991798, 993767], Y=[211358, 213327]",
        "Ground elevation: 5.51 ft MSL (1.7 m)",
        "",
        "Tallest building: 169 m (BIN 1090642, built 2018)",
        "UN Secretariat: 156 m (BIN 1083875)",
        "One UN Plaza: 154 m (BIN 1038758)",
    ], size=10)

    # Right column: results
    textbox(slide, 5.2, 0.7, 4.5, 0.35, "Results", size=14, bold=True, color=NAVY)
    bullet_list(slide, 5.4, 1.05, 4.3, 2.5, [
        "CityGML (primary): 123 buildings (93%)",
        "Socrata fallback: 8 buildings (6%)",
        "Total: 131 buildings",
        "Total triangles: 23,575",
        "",
        "Surface breakdown:",
        "  GroundSurface: 125",
        "  RoofSurface: 1,048",
        "  WallSurface: 6,211",
        "",
        "Bounding box: 732 m × 731 m × 170 m",
    ], size=10)

    # Bottom: output
    textbox(slide, 0.4, 3.8, 9.2, 0.35, "Output", size=14, bold=True, color=NAVY)
    bullet_list(slide, 0.6, 4.15, 9, 1, [
        "File: un_hq_surroundings.stl (binary STL, 23,575 triangles)",
        "Coordinate system: local meters centered at UN Secretariat, Z=0 at ground",
        "Solver-agnostic: compatible with FEFLO, OpenFOAM, Fluent, or any STL-reading CFD code",
    ], size=10)


def add_methodology(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_title(slide, "Methodology: Hybrid CityGML + Socrata Pipeline")
    bullet_list(slide, 0.4, 0.7, 9.2, 4.5, [
        "STEP 1: CITYGML EXTRACTION (primary source — 2014 aerial photogrammetry)",
        "  1a. Convert domain corners from WGS84 to EPSG:2263 (NY State Plane feet)",
        "  1b. Identify overlapping DA district files (DA12 covers Midtown + Turtle Bay)",
        "  1c. Stream through GML files line-by-line (764 MB for DA12, no full-file load)",
        "  1d. For each <bldg:Building>: check if any coordinate falls in domain bbox",
        "  1e. Parse WallSurface, RoofSurface, GroundSurface → fan triangulate",
        "  1f. Convert EPSG:2263 feet → local meters (centered at target building)",
        "  1g. Auto-detect ground elevation from GroundSurface Z values",
        "",
        "STEP 2: SOCRATA API FALLBACK (for post-2014 construction)",
        "  2a. Query NYC Building Footprints API (GeoJSON) within domain bbox",
        "  2b. Deduplicate: skip if BIN matches CityGML set (71 caught)",
        "  2c. Deduplicate: skip if centroid within 15m of any CityGML centroid",
        "  2d. Extrude remaining footprint polygons to height_roof (LoD1, flat-topped)",
        "  2e. Result: 8 genuinely new buildings added",
        "",
        "STEP 3: OUTPUT",
        "  3a. Merge all triangles → un_hq_surroundings.stl (binary STL)",
        "  3b. Print provenance report with building counts + triangle counts",
    ], size=10)


def add_data_sources(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_title(slide, "Data Sources")

    textbox(slide, 0.4, 0.7, 4.3, 0.35, "NYC 3D Building Model (CityGML)", size=14, bold=True, color=NAVY)
    bullet_list(slide, 0.6, 1.05, 4.1, 2.5, [
        "2014 aerial photogrammetry survey",
        "~1 million buildings across 5 boroughs",
        "LoD2: actual wall, roof, ground surfaces",
        "12 DA district files (~6 GB total)",
        "EPSG:2263 coordinate system (US Survey Feet)",
        "Source: georocket/new-york-city-model-enhanced",
        "GitHub release: 20v5 (enhanced version)",
    ], size=10)

    textbox(slide, 5.2, 0.7, 4.3, 0.35, "Socrata Building Footprints API", size=14, bold=True, color=NAVY)
    bullet_list(slide, 5.4, 1.05, 4.1, 2.5, [
        "Continuously updated (includes 2015+ construction)",
        "2D footprint polygon + height_roof attribute",
        "BIN (Building Identification Number) for matching",
        "GeoJSON format, WGS84 coordinates",
        "Extruded to flat-topped LoD1 boxes",
        "API endpoint: data.cityofnewyork.us/resource/5zhs-2jue",
    ], size=10)

    textbox(slide, 0.4, 3.5, 9.2, 0.35, "Deduplication Strategy", size=14, bold=True, color=NAVY)
    bullet_list(slide, 0.6, 3.85, 9, 1.5, [
        "Step 1: BIN matching — if Socrata BIN exists in CityGML set, skip (caught 71)",
        "Step 2: Spatial proximity — if Socrata centroid within 15m of CityGML centroid, skip",
        "Step 3: Remainder kept as fallback — genuinely new post-2014 buildings (8 added)",
    ], size=10)


def add_reproduction(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_title(slide, "How to Reproduce")

    bullet_list(slide, 0.4, 0.7, 9.2, 4.5, [
        "PREREQUISITES",
        "  Python 3.8+, requests library (pip install requests)",
        "  CityGML data: _citygml/DA12.gml.zip (52 MB, auto-downloaded from GitHub)",
        "",
        "GENERATE STL (single command):",
        "  python generate_stl_hybrid.py \\",
        "      --center-lat 40.7489 --center-lon -73.9680 \\",
        "      --name un_hq --radius 300 --no-target",
        "",
        "  Output: constant/triSurface/un_hq_surroundings.stl",
        "  Runtime: ~2 minutes (DA12 scan: 20s, Socrata API: 2s)",
        "",
        "OPTIONAL FLAGS:",
        "  --offline          Skip Socrata API (CityGML only)",
        "  --year 1978        Historical filter (buildings built by year)",
        "  --ascii            Also write ASCII STL for inspection",
        "  --data-dir PATH    Custom CityGML data directory",
        "  --output-dir PATH  Custom STL output directory",
        "",
        "VISUALIZE:",
        "  pvpython _paraview_un_hq.py   (generates 5 PNG views)",
        "  Open building_browser.html → click 'UN Headquarters' preset",
    ], size=10)


def add_building_inventory(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_title(slide, "Notable Buildings in Domain")

    headers = ["Building", "Address", "Height", "Source"]
    rows = [
        ["Turkish Consulate (2018)", "821 1 Avenue", "169 m", "Socrata"],
        ["50 UN Plaza (2016)", "50 United Nations Plaza", "163 m", "Socrata"],
        ["UN Secretariat", "752 1 Avenue", "156 m", "CityGML"],
        ["One UN Plaza", "787 1 Avenue", "154 m", "CityGML"],
        ["UN Dev. Program", "323 E 44 Street", "147 m", "CityGML"],
        ["685 First Ave (2018)", "685 First Avenue", "140 m", "Socrata"],
        ["US Mission to UN", "799 1 Avenue", "117 m", "CityGML"],
        ["Woodstock Tower", "312 E 42 Street", "111 m", "CityGML"],
        ["Nigerian Mission", "828 2 Avenue", "94 m", "CityGML"],
        ["Windsor Tower", "1 Tudor City Place", "89 m", "CityGML"],
        ["Ford Foundation", "320 E 43 Street", "47 m", "CityGML"],
    ]

    col_widths = [2.3, 2.3, 0.8, 1.0]
    col_x = [0.3, 2.7, 5.1, 6.1]

    y = 0.7
    for j, h in enumerate(headers):
        textbox(slide, col_x[j], y, col_widths[j], 0.3, h, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        tb = slide.shapes[-1]
        tb.fill.solid()
        tb.fill.fore_color.rgb = NAVY

    y += 0.35
    for row in rows:
        for j, cell in enumerate(row):
            c = GREEN if cell == "Socrata" else DARK
            b = cell == "Socrata"
            textbox(slide, col_x[j], y, col_widths[j], 0.22, cell, size=9, color=c, bold=b,
                    align=PP_ALIGN.CENTER if j > 1 else PP_ALIGN.LEFT)
        y += 0.23

    bullet_list(slide, 0.5, y + 0.2, 9, 0.8, [
        "131 total buildings | 123 from CityGML LoD2 (actual surfaces) | 8 from Socrata (flat-topped)",
        "CityGML buildings have roof shapes, wall setbacks from 2014 aerial photogrammetry",
        "Socrata buildings are post-2014 construction extruded to flat tops (LoD1)",
    ], size=9, color=GRAY)


def add_screenshot_slide(prs, img_file, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_title(slide, title)

    img_path = os.path.join(PARENT_DIR, img_file)
    if os.path.exists(img_path):
        # Full-width centered image
        slide.shapes.add_picture(img_path, Inches(0.3), Inches(0.65),
                                 Inches(9.4), Inches(4.5))
    else:
        textbox(slide, 2, 2.5, 6, 1, f"Screenshot not found: {img_file}",
                size=14, color=GRAY, align=PP_ALIGN.CENTER)


def add_pipeline_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_title(slide, "Pipeline Architecture")

    bullet_list(slide, 0.4, 0.7, 4.5, 4, [
        "STREAMING PARSER",
        "  Line-by-line scan of GML files (764 MB+)",
        "  O(1) memory: only current building in RAM",
        "  Bounding box check on first coordinate → skip early",
        "  Pattern: <bldg:Building> ... </bldg:Building>",
        "",
        "COORDINATE TRANSFORMS",
        "  CityGML: EPSG:2263 (US Survey Feet)",
        "    → subtract center → multiply by 1200/3937",
        "  Socrata: WGS84 (degrees)",
        "    → same scale factors for alignment",
        "  Ground: auto-detected from GroundSurface Z",
        "",
        "TRIANGULATION",
        "  CityGML polygons: fan triangulation from v0",
        "  Socrata footprints: ear-clipping + extrusion",
    ], size=10)

    bullet_list(slide, 5.2, 0.7, 4.5, 4, [
        "LOCATION-AGNOSTIC DESIGN",
        "  --center-lat/--center-lon for any NYC building",
        "  --radius for domain size",
        "  --no-target for surroundings-only mode",
        "  --exclude-bins to remove target building",
        "",
        "DA DISTRICT AUTO-DETECTION",
        "  20 pre-computed EPSG:2263 envelopes",
        "  Domain → overlapping districts → scan only those",
        "  Auto-downloads missing districts from GitHub",
        "",
        "OUTPUT",
        "  Binary STL (watertight, outward normals)",
        "  Compatible with any CFD solver",
        "  FEFLO, OpenFOAM, Fluent, Star-CCM+",
        "  Optional ASCII STL for inspection",
    ], size=10)


def add_summary_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY
    textbox(slide, 0.5, 1.2, 9, 1,
            "UN HQ Area — Ready for CFD", size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    bullet_list(slide, 1, 2.6, 8, 2.5, [
        "131 buildings extracted (CityGML LoD2 + Socrata fallback)",
        "23,575 triangles with actual roof shapes and wall setbacks",
        "Single command: python generate_stl_hybrid.py --center-lat 40.7489 ...",
        "Pipeline works for any NYC location (Manhattan, Brooklyn, etc.)",
        "Building browser (HTML) for visual QA against satellite imagery",
    ], size=14, color=WHITE)


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    add_title_slide(prs)             # 1
    add_parameters(prs)              # 2
    add_methodology(prs)             # 3
    add_data_sources(prs)            # 4
    add_reproduction(prs)            # 5
    add_pipeline_architecture(prs)   # 6
    add_building_inventory(prs)      # 7

    for img, title in SCREENSHOTS:   # 8-12
        add_screenshot_slide(prs, img, f"UN HQ — {title}")

    add_summary_slide(prs)           # 13

    prs.save(OUT_PATH)
    print(f"Saved: {OUT_PATH} ({len(prs.slides)} slides)")


if __name__ == '__main__':
    main()
