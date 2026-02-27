#!/usr/bin/env python3
"""Create PowerPoint presentation: NYC Mesh Methods & Limitations.

Explains the data pipeline for creating accurate urban building meshes
for CFD simulation of the Citicorp Center wind loading case.

Usage: python _create_mesh_presentation.py
Output: mesh_methods_presentation.pptx (in parent directory)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

# Colors
DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
MED_BLUE = RGBColor(0x2C, 0x5F, 0x8A)
LIGHT_BLUE = RGBColor(0xD6, 0xE8, 0xF5)
ACCENT_GREEN = RGBColor(0x2E, 0x7D, 0x32)
ACCENT_RED = RGBColor(0xC6, 0x28, 0x28)
ACCENT_ORANGE = RGBColor(0xE6, 0x77, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)


def set_slide_bg(slide, color):
    """Set slide background color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=BLACK, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    """Add a text box to the slide."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_slide(slide, items, left=0.8, top=2.0, width=8.5, font_size=16,
                     color=BLACK, spacing=Pt(8)):
    """Add bulleted text to a slide."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(5.0))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # Handle indented sub-bullets
        if isinstance(item, tuple):
            text, level = item
            p.level = level
        else:
            text = item
            p.level = 0

        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.space_after = spacing

    return txBox


def add_table(slide, rows, cols, data, left, top, width, height,
              header_color=DARK_BLUE, header_font_color=WHITE):
    """Add a table to the slide."""
    table_shape = slide.shapes.add_table(rows, cols,
                                          Inches(left), Inches(top),
                                          Inches(width), Inches(height))
    table = table_shape.table

    for row_idx in range(rows):
        for col_idx in range(cols):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(data[row_idx][col_idx])

            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(12)
                paragraph.font.name = 'Calibri'
                if row_idx == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = header_font_color
                else:
                    paragraph.font.color.rgb = BLACK

            if row_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color

    return table


def add_title_bar(slide, title, subtitle=None):
    """Add a colored title bar at the top."""
    # Title background
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(0), Inches(0),
                                    Inches(10), Inches(1.3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()

    # Title text
    add_text_box(slide, 0.5, 0.15, 9, 0.6, title,
                 font_size=28, bold=True, color=WHITE)

    if subtitle:
        add_text_box(slide, 0.5, 0.75, 9, 0.5, subtitle,
                     font_size=14, color=LIGHT_BLUE)


def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ============================================================
    # SLIDE 1: Title
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide, WHITE)

    # Big title bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(0), Inches(0),
                                    Inches(10), Inches(3.0))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()

    add_text_box(slide, 0.8, 0.6, 8.5, 1.0,
                 "Accurate NYC Building Meshes\nfor CFD Simulation",
                 font_size=36, bold=True, color=WHITE)

    add_text_box(slide, 0.8, 2.0, 8.5, 0.5,
                 "Citicorp Center Wind Loading Analysis",
                 font_size=20, color=LIGHT_BLUE)

    add_text_box(slide, 0.8, 3.5, 8.5, 0.5,
                 "OR 750 -- Reliability, Safety, and Risk",
                 font_size=18, color=GRAY)

    add_text_box(slide, 0.8, 4.2, 8.5, 0.5,
                 "Data Sources  |  Methods  |  Validation",
                 font_size=16, color=GRAY)

    # ============================================================
    # SLIDE 2: Problem Statement
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Why Accurate Building Geometry Matters",
                  "CFD wind simulation requires realistic urban context")

    add_bullet_slide(slide, [
        "Wind loads on the Citicorp Center depend on surrounding buildings",
        ("Upstream buildings create turbulence and wake effects", 1),
        ("Building blockage alters pressure distributions on target", 1),
        ("Missing or inaccurate neighbors = wrong wind loads", 1),
        "",
        "Our original mesh used 18 hand-estimated box buildings",
        ("Approximate positions from Google Maps visual inspection", 1),
        ("Uniform flat-top extrusions (no roof shapes)", 1),
        ("Only ~276 triangles for surroundings geometry", 1),
        "",
        "Goal: Replace with surveyed NYC building data",
    ], top=1.6, font_size=17)

    # ============================================================
    # SLIDE 3: Data Sources Overview
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "NYC Building Data Sources",
                  "Three levels of detail available")

    data = [
        ["Source", "Date", "Detail", "Format", "Status"],
        ["NYC Building Footprints\n(Socrata API)", "Continuously\nupdated",
         "LoD1: footprint +\nroof height", "GeoJSON\n(2D polygon + height)",
         "Used as\nfallback"],
        ["NYC 3D Building Model\n(CityGML)", "Aug 2014\naerial survey",
         "LoD1-2: actual\nroof/wall surfaces", "CityGML XML\n(3D polygons)",
         "Primary\nsource"],
        ["NYC 2017 Lidar\n(Topobathymetric)", "May 2017",
         "Raw point cloud\n(8 pts/m^2)", "LAS 1.4\n(classified)", "Available\n(not used)"],
    ]
    add_table(slide, 4, 5, data, 0.5, 1.6, 9.0, 3.5)

    add_text_box(slide, 0.5, 5.5, 9.0, 1.5,
                 "Decision: CityGML provides the best balance of accuracy and usability.\n"
                 "Actual roof shapes and wall setbacks without point cloud processing.\n"
                 "Socrata API fills the gap for buildings constructed after 2014.",
                 font_size=14, color=GRAY)

    # ============================================================
    # SLIDE 4: LoD1 vs LoD2 Explanation
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "LoD1 vs LoD2: What's the Difference?",
                  "Level of Detail in 3D city models")

    # LoD1 column
    add_text_box(slide, 0.5, 1.6, 4.0, 0.5,
                 "LoD1 -- Extruded Footprint", font_size=20, bold=True, color=ACCENT_RED)
    add_bullet_slide(slide, [
        "2D polygon extruded to a single height",
        "Flat roof (no shape detail)",
        "No setbacks or terracing",
        "Good for distant buildings",
        "Source: Socrata API (height_roof field)",
    ], left=0.5, top=2.3, width=4.0, font_size=14)

    # LoD2 column
    add_text_box(slide, 5.5, 1.6, 4.0, 0.5,
                 "LoD2 -- Surface Model", font_size=20, bold=True, color=ACCENT_GREEN)
    add_bullet_slide(slide, [
        "Actual WallSurface polygons at multiple Z",
        "RoofSurface with slopes and ridges",
        "GroundSurface at surveyed elevation",
        "Captures setbacks, terraces, crowns",
        "Source: CityGML (2014 aerial photogrammetry)",
    ], left=5.5, top=2.3, width=4.0, font_size=14)

    # Divider line
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(4.85), Inches(1.6),
                                    Inches(0.03), Inches(4.0))
    shape.fill.solid()
    shape.fill.fore_color.rgb = GRAY
    shape.line.fill.background()

    # Bottom comparison
    add_text_box(slide, 0.5, 5.8, 9.0, 1.0,
                 "Example: Citicorp tower -- LoD1 = 12 triangles (flat box), "
                 "LoD2 = 124 triangles (17 Z levels, irregular crown).\n"
                 "For surrounding buildings, LoD1 is acceptable (blockage dominates). "
                 "For the target building, LoD2 provides aerodynamically meaningful detail.",
                 font_size=13, color=GRAY)

    # ============================================================
    # SLIDE 5: Hybrid Pipeline
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Hybrid Data Pipeline",
                  "CityGML primary + Socrata API fallback")

    add_bullet_slide(slide, [
        "Step 1: Stream through CityGML DA district files (~764 MB each)",
        ("Line-by-line parser to avoid loading into memory", 1),
        ("Extract buildings whose coordinates fall in CFD domain", 1),
        ("Parse WallSurface, RoofSurface, GroundSurface polygons", 1),
        "",
        "Step 2: Query Socrata API for buildings in same domain",
        ("Returns current building inventory with heights", 1),
        ("Match by BIN (Building Identification Number)", 1),
        "",
        "Step 3: Merge -- CityGML geometry where available, Socrata extrusion for rest",
        ("Buildings in both sources: use CityGML (more accurate)", 1),
        ("Buildings only in Socrata (built after 2014): extrude footprint", 1),
        ("Citicorp tower + stilts: always use hand-crafted geometry", 1),
    ], top=1.6, font_size=15)

    # ============================================================
    # SLIDE 6: Coordinate Systems
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Coordinate Systems & Transforms",
                  "Three coordinate systems in the pipeline")

    data = [
        ["System", "Units", "Source", "Transform"],
        ["WGS84 (EPSG:4326)", "Degrees (lat/lon)", "Socrata API,\nuser input",
         "Equirectangular\nprojection"],
        ["NY State Plane\n(EPSG:2263)", "US Survey Feet", "CityGML\nbuilding data",
         "Linear scaling\n(1200/3937 m/ft)"],
        ["Local Cartesian", "Meters", "STL output\n(OpenFOAM)", "Origin at Citicorp\nZ=0 at ground"],
    ]
    add_table(slide, 4, 4, data, 0.5, 1.6, 9.0, 2.5)

    add_text_box(slide, 0.5, 4.5, 9.0, 2.5,
                 "Key parameters:\n"
                 "  Citicorp Center: 40.7579N, 73.9690W\n"
                 "  EPSG:2263 reference: X = 992,506 ft, Y = 215,620 ft\n"
                 "  Ground elevation: 25.78 ft MSL (from CityGML GroundSurface)\n"
                 "  US Survey Foot: 1200/3937 m = 0.30480061 m (NOT 0.3048)\n\n"
                 "Linear approximation accurate to ~1m within 10km of reference point.\n"
                 "No pyproj or GDAL dependencies needed.",
                 font_size=13, color=GRAY, font_name='Consolas')

    # ============================================================
    # SLIDE 7: Results Comparison
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Results: Old vs. Hybrid Mesh",
                  "Quantitative comparison of surrounding building geometry")

    data = [
        ["Metric", "Old (Socrata LoD1)", "Hybrid (CityGML + Socrata)", "Improvement"],
        ["Buildings", "~170 (API)", "806 (599 CityGML + 207 Socrata)", "4.7x"],
        ["Triangles", "~14,000", "64,486", "4.6x"],
        ["Surface types", "Flat-top extrusion only", "Wall + Roof + Ground surfaces", "LoD2"],
        ["Roof detail", "None (all flat)", "Actual shapes for 599 buildings", "Realistic"],
        ["Data currency", "Present-day only", "2014 CityGML + present-day fill", "Complete"],
        ["Coverage", "Socrata: 74%\nof hybrid total", "CityGML: 74%\nSocrata fallback: 26%", "--"],
    ]
    add_table(slide, 7, 4, data, 0.3, 1.6, 9.4, 4.0)

    add_text_box(slide, 0.5, 6.0, 9.0, 1.0,
                 "BIN matching: 499 Socrata buildings matched to CityGML and replaced with "
                 "higher-quality geometry.\n207 buildings only in Socrata (built after 2014) -- "
                 "extruded as LoD1 fallback.\n3 Citicorp complex BINs excluded (hand-crafted tower + stilts).",
                 font_size=13, color=GRAY)

    # ============================================================
    # SLIDE 8: Limitations of Old Mesh
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Limitations of the Original Mesh",
                  "Why the hand-estimated geometry was insufficient")

    add_bullet_slide(slide, [
        "Only 18 buildings -- Midtown Manhattan has 800+ in the CFD domain",
        ("Missing buildings = missing blockage = wrong wake patterns", 1),
        "",
        "Positions estimated from Google Maps satellite view",
        ("Typical error: 10-50m in X/Y placement", 1),
        ("No ground truth for building heights", 1),
        "",
        "All buildings were axis-aligned rectangular boxes",
        ("Real footprints are L-shaped, irregular, with setbacks", 1),
        ("All roofs flat -- real roofs have slopes and crowns", 1),
        "",
        "No temporal filtering capability",
        ("Could not distinguish 1978 vs. 2024 skyline", 1),
        ("Critical for studying the original structural crisis", 1),
        "",
        "276 total triangles vs. 64,486 in the hybrid mesh",
    ], top=1.6, font_size=16)

    # ============================================================
    # SLIDE 9: Citicorp Tower Detail
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Citicorp Tower: Simplified vs. LoD2",
                  "BIN 1036474 -- Individual Landmark, 60 floors, built 1978")

    data = [
        ["Property", "Simplified STL", "CityGML LoD2", "Difference"],
        ["Triangles", "19 (12 box + 7 crown)", "124", "+105 (6.5x)"],
        ["X/Y span", "67.7 m (axis-aligned)", "64.4 m (true outline)", "-3.3 m"],
        ["Height", "244.2 m (stilt to crown)", "283.7 m (ground to crown)", "+39.6 m"],
        ["Z levels", "4", "17", "+13 levels"],
        ["Cross-section", "Perfect 45-degree diamond", "Irregular polygon\nwith setbacks", "More realistic"],
        ["Crown", "3-level approximation", "~10 distinct Z levels\n(241-284 m)", "Detailed terracing"],
        ["Stilts", "4 axis-aligned boxes", "Not in CityGML\n(hidden from aerial)", "Keep hand-crafted"],
    ]
    add_table(slide, 8, 4, data, 0.3, 1.5, 9.4, 4.5)

    add_text_box(slide, 0.5, 6.3, 9.0, 0.8,
                 "Note: CityGML height (283.7m) exceeds actual 279m because rooftop "
                 "equipment was captured by the photogrammetric survey.\n"
                 "Stilts are not visible in aerial imagery -- hand-crafted geometry "
                 "is retained for all configurations.",
                 font_size=13, color=GRAY)

    # ============================================================
    # SLIDE 10: Post-2014 Gap
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Filling the Post-2014 Gap",
                  "Socrata fallback captures buildings missing from CityGML")

    add_bullet_slide(slide, [
        "NYC 3D Building Model is from 2014 aerial survey -- never updated",
        "",
        "Notable buildings missing from CityGML near Citicorp:",
    ], top=1.5, font_size=16)

    data = [
        ["Building", "Distance", "Height", "Completed", "CFD Impact"],
        ["Selene (100 E 53rd)", "~50 m", "217 m (63 fl)", "2019", "HIGH"],
        ["270 Park Ave\n(JPMorgan HQ)", "~600 m", "423 m (60 fl)", "2025", "MEDIUM"],
        ["One Vanderbilt", "~1.2 km", "427 m (62 fl)", "2020", "LOW"],
    ]
    add_table(slide, 4, 5, data, 0.5, 3.3, 9.0, 2.0)

    add_text_box(slide, 0.5, 5.7, 9.0, 1.5,
                 "Selene is the most critical missing building: at 217m and catty-corner "
                 "to Citicorp, it significantly alters the local wind field.\n\n"
                 "The Socrata API provides current building inventory. "
                 "207 post-2014 buildings were added via the fallback pipeline, "
                 "including all buildings listed above.",
                 font_size=14, color=GRAY)

    # ============================================================
    # SLIDE 11: Building Extraction Tool
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "General-Purpose Extraction Tool",
                  "extract_buildings.py -- extract any NYC location")

    add_text_box(slide, 0.5, 1.6, 9.0, 0.5,
                 "python extract_buildings.py --lat 40.7580 --lon -73.9690 --radius 300",
                 font_size=16, bold=True, color=MED_BLUE, font_name='Consolas')

    add_bullet_slide(slide, [
        "Self-contained script: Python 3.8+ stdlib only (no pip installs)",
        "Auto-detects DA district files, downloads from GitHub if missing",
        "Streams through large GML files (~764 MB) line by line",
        "Outputs: CityGML fragments + binary STL (OpenFOAM-ready)",
        "",
        "Attributes extracted per building:",
        ("BIN (Building Identification Number)", 1),
        ("Year built, number of floors, landmark status", 1),
        ("Height from roof/ground surface geometry", 1),
        "",
        "Tested: 31 buildings at Citicorp (100m radius), 4,392 triangles, ~40 seconds",
    ], top=2.3, font_size=15)

    # ============================================================
    # SLIDE 12: Technical Summary
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Technical Summary",
                  "What the hybrid pipeline provides")

    # Left column
    add_text_box(slide, 0.5, 1.5, 4.3, 0.4,
                 "Capabilities", font_size=18, bold=True, color=ACCENT_GREEN)
    add_bullet_slide(slide, [
        "806 buildings in CFD domain",
        "64,486 triangles (4.6x old mesh)",
        "LoD2 surfaces for 74% of buildings",
        "Present-day + historical (--year 1978)",
        "BIN-based deduplication",
        "Offline mode (CityGML only)",
        "Hardcoded fallback (both fail)",
        "Binary STL (OpenFOAM-ready)",
    ], left=0.5, top=2.1, width=4.3, font_size=14, color=BLACK)

    # Right column
    add_text_box(slide, 5.3, 1.5, 4.3, 0.4,
                 "Limitations", font_size=18, bold=True, color=ACCENT_RED)
    add_bullet_slide(slide, [
        "CityGML from 2014 (12 years old)",
        "LoD2 detail only for ~100 iconic bldgs",
        "Most CityGML buildings are LoD1",
        "Fan triangulation (not ear-clipping)",
        "No terrain elevation variation",
        "Stilts not in aerial data",
        "DA envelope matching too broad",
        "Linear coord transform (~0.1% error)",
    ], left=5.3, top=2.1, width=4.3, font_size=14, color=BLACK)

    # Divider
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(4.95), Inches(1.5),
                                    Inches(0.03), Inches(4.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = GRAY
    shape.line.fill.background()

    # ============================================================
    # SLIDE 13: Next Steps
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Next Steps")

    add_bullet_slide(slide, [
        "1. Run snappyHexMesh with hybrid geometry (target: 8-10M cells)",
        "",
        "2. Compare drag coefficients:",
        ("Old mesh (18 hand-estimated buildings) vs. hybrid (806 surveyed)", 1),
        ("Present-day skyline vs. 1978 skyline (--year 1978 filter)", 1),
        "",
        "3. Correct stilt geometry",
        ("Rotate 45 degrees to match tower orientation", 1),
        ("Verify placement at face midpoints", 1),
        "",
        "4. Validate against published wind tunnel data",
        ("Camelli et al. (2006) -- VLES of NYC wind fields", 1),
        ("ASCE 7-22 wind load provisions", 1),
        "",
        "5. ParaView visualization of building geometry comparisons",
    ], top=1.6, font_size=16)

    # ============================================================
    # SAVE
    # ============================================================
    script_dir = os.path.dirname(os.path.abspath(__file__))
    parent_dir = os.path.dirname(script_dir)
    out_path = os.path.join(parent_dir, 'mesh_methods_presentation.pptx')
    prs.save(out_path)
    print(f"Saved: {out_path}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == '__main__':
    create_presentation()
